import streamlit as st
import re
import requests
import json
from pptx import Presentation
import openai
import time

# HTML 태그를 제거하는 함수
def clean_html_tags(text):
    # 정규 표현식으로 HTML 태그를 제거
    clean_text = re.sub(r'<.*?>', '', text)
    return clean_text

# 네이버 맞춤법 검사 API 클래스
class SpellChecker:
    def __init__(self, api_key=None, use_naver_api=False):
        self.api_key = api_key
        self.use_naver_api = use_naver_api
        self.passport_key = None
        self.base_url = None

    # 네이버 맞춤법 검사 API에서 필요한 passport_key를 가져오는 함수
    def fetch_passport_key(self):
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/129.0.0.0 Safari/537.36',
            'Referer': 'https://search.naver.com/',
        }
        response = requests.get("https://search.naver.com/search.naver?query=%EB%A7%9E%EC%B6%A9%EB%B2%95%20%EA%B2%80%EC%82%AC%EA%B8%B0", headers=headers)
        passport_key_match = re.search(r'(?<={new SpellingCheck\({API:{checker:").*?(?="},selector)', response.text)
        if not passport_key_match:
            raise Exception("Error: Unable to retrieve passport key")
        self.base_url, self.passport_key = passport_key_match.group(0).split("?passportKey=")

    # 네이버 맞춤법 검사 API 호출
    def spell_check_with_naver(self, text):
        if self.passport_key is None or self.base_url is None:
            self.fetch_passport_key()
        payload = {
            'passportKey': self.passport_key,
            'where': 'nexearch',
            'color_blindness': 0,
            'q': text
        }
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            'Referer': 'https://search.naver.com/',
        }
        result_response = requests.get(self.base_url, headers=headers, params=payload)
        return json.loads(result_response.text)['message']['result']['notag_html']

    # OpenAI API로 맞춤법 검사
    def spell_check_with_openai(self, text):
        openai.api_key = self.api_key
        
        prompt = f"다음 텍스트의 맞춤법을 단어별로 교정해주세요.:\n{text}. 교정한 단어 별로 띄어쓰기, 오탈자, 문법오류, 조사오류, 문장부호 오류 중 어느 오류인지도 말하세요."
        sys_prmt = "너는 아주 차가운 맞춤법 교정 AI야. 필요한 단어만 말하고 쓸떼없는 말은 안 해. 교정할 거 없으면 하지마. 외래어 굳이 순화하지마."
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt},{"role": "system", "content": sys_prmt}],
            max_tokens=500,
            temperature=0.8
        )
        return response['choices'][0]['message']['content'].strip()


    # 맞춤법 검사 함수 (네이버 또는 OpenAI API 호출)
    def spell_check(self, text):
        if self.use_naver_api:
            return self.spell_check_with_naver(text)
        else:
            return self.spell_check_with_openai(text)

# PPT 파일에서 텍스트 추출하는 함수
def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        slides_text.append("\n".join(slide_text))
    return slides_text

# PPT 파일과 맞춤법 검사기를 사용하여 결과 반환
def process_ppt_with_spellchecker(ppt_file, checker):
    extracted_text = extract_text_from_ppt(ppt_file)
    results = []
    for page, text in enumerate(extracted_text, start=1):
        corrected = checker.spell_check(text)
        results.append((page, text, corrected))
    return results

# Streamlit 앱 실행 함수
def run(api_key=None):
    st.title("PPT 맞춤법 검사")

    # 사용자에게 API 선택을 위한 옵션 제공
    api_option = st.radio(
        "사용할 맞춤법 검사 API 선택",
        ('OpenAI API', '네이버 맞춤법 검사 API'),
        index=0  # 기본값은 OpenAI API로 설정
    )

    ppt_file = st.file_uploader("PPT 파일을 업로드하세요.", type="pptx")
    analyze_button = st.button("맞춤법 검사")

    if ppt_file and analyze_button:
        use_naver_api = api_option == '네이버 맞춤법 검사 API'  # 사용자가 선택한 API에 따라 설정
        checker = SpellChecker(api_key=api_key, use_naver_api=use_naver_api)  # SpellChecker 객체 초기화
        
        # 로딩 중 표시
        with st.spinner('맞춤법 검사 중...'):
            # ppt 파일에 대해 맞춤법 검사 실행
            spell_check_results = process_ppt_with_spellchecker(ppt_file, checker)
            time.sleep(1)  # 로딩이 진행되는 것처럼 잠시 대기 (옵션)
        
        # 결과를 페이지별로 토글 형태로 표시
        st.subheader("분석 결과")
        
        for page, original, corrected in spell_check_results:
            # HTML 태그를 제거한 텍스트를 사용하여 결과 출력
            cleaned_corrected = clean_html_tags(corrected)
            
            # 네이버 API인 경우 수정 전후 비교
            if use_naver_api:
                with st.expander(f"{page}페이지", expanded=True):
                    st.markdown(f"▪️ {original}➡️ {cleaned_corrected}")
            # OpenAI API인 경우 수정된 결과만 표시
            else:
                with st.expander(f"{page}페이지", expanded=True):
                    st.write(cleaned_corrected)
        
        # 사용한 API 표시
        if use_naver_api:
            st.markdown("**네이버 맞춤법 검사 API를 사용했습니다.**")
        else:
            st.markdown("**OpenAI API를 사용했습니다.**")
